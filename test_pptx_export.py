#!/usr/bin/env python3
"""Test script for PPTX export functionality - bypassing MCP."""

import asyncio
import json
import sys
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional
import io

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

# Export directory configuration
EXPORT_DIR = "/tmp/protex-intelligence-file-exports"

# Hardcoded test data
TEST_SLIDES = [
    {
        "backgroundColor": "F8F9FA",
        "elements": [
            {
                "type": "text",
                "text": "Q4 2024 Performance Report",
                "options": {
                    "x": 1, "y": 0.5, "w": 11, "h": 1.5,
                    "fontSize": 36, "color": "2C3E50", "bold": True
                }
            },
            {
                "type": "text", 
                "text": "Comprehensive analysis of company metrics and growth",
                "options": {
                    "x": 1, "y": 2, "w": 11, "h": 1,
                    "fontSize": 18, "color": "7F8C8D"
                }
            },
            {
                "type": "shape",
                "shapeType": "rectangle",
                "options": {
                    "x": 1, "y": 5.5, "w": 11, "h": 0.5,
                    "fill": "3498DB"
                }
            }
        ]
    },
    {
        "backgroundColor": "FFFFFF",
        "elements": [
            {
                "type": "text",
                "text": "Sales Performance",
                "options": {
                    "x": 1, "y": 0.5, "w": 6, "h": 1,
                    "fontSize": 28, "color": "2C3E50", "bold": True
                }
            },
            {
                "type": "chart",
                "chartType": "bar",
                "chartData": [
                    {
                        "name": "Q4 Sales",
                        "labels": ["Oct", "Nov", "Dec"],
                        "values": [850000, 920000, 1050000]
                    },
                    {
                        "name": "Q3 Sales",
                        "labels": ["Oct", "Nov", "Dec"], 
                        "values": [780000, 820000, 890000]
                    }
                ],
                "options": {
                    "x": 1, "y": 1.5, "w": 11, "h": 4,
                    "title": "Monthly Sales Comparison"
                }
            }
        ]
    },
    {
        "backgroundColor": "FFFFFF",
        "elements": [
            {
                "type": "text",
                "text": "Team Performance Metrics",
                "options": {
                    "x": 1, "y": 0.5, "w": 11, "h": 1,
                    "fontSize": 28, "color": "2C3E50", "bold": True
                }
            },
            {
                "type": "table",
                "rows": [
                    ["Department", "Target", "Actual", "Growth %"],
                    ["Sales", "$2.5M", "$2.82M", "+12.8%"],
                    ["Marketing", "500 leads", "645 leads", "+29.0%"],
                    ["Support", "90% satisfaction", "94% satisfaction", "+4.4%"],
                    ["Development", "15 features", "18 features", "+20.0%"]
                ],
                "options": {
                    "x": 1, "y": 2, "w": 11, "h": 3
                }
            }
        ]
    },
    {
        "backgroundColor": "FFFFFF", 
        "elements": [
            {
                "type": "text",
                "text": "Key Achievements",
                "options": {
                    "x": 1, "y": 0.5, "w": 11, "h": 1,
                    "fontSize": 28, "color": "2C3E50", "bold": True
                }
            },
            {
                "type": "text",
                "text": [
                    {"text": "✓ ", "color": "27AE60", "fontSize": 20, "bold": True},
                    {"text": "Exceeded sales targets by 12.8%\n", "fontSize": 18},
                    {"text": "✓ ", "color": "27AE60", "fontSize": 20, "bold": True},
                    {"text": "Increased customer satisfaction to 94%\n", "fontSize": 18},
                    {"text": "✓ ", "color": "27AE60", "fontSize": 20, "bold": True},
                    {"text": "Delivered 20% more features than planned\n", "fontSize": 18},
                    {"text": "✓ ", "color": "27AE60", "fontSize": 20, "bold": True},
                    {"text": "Generated 29% more marketing leads", "fontSize": 18}
                ],
                "options": {
                    "x": 1, "y": 2, "w": 11, "h": 4
                }
            }
        ]
    }
]

TEST_ARGUMENTS = {
    "slides": TEST_SLIDES,
    "filename": "quarterly_performance_report",
    "description": "Q4 2024 comprehensive performance analysis presentation",
    "options": {
        "layout": "16x9",
        "author": "Business Analytics Team",
        "title": "Q4 2024 Performance Report",
        "subject": "Quarterly Business Review"
    }
}

# Import the functions from server.py
def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor object."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )

def get_chart_type(chart_type_str: str) -> int:
    """Map chart type string to python-pptx chart type."""
    chart_type_map = {
        'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'line': XL_CHART_TYPE.LINE,
        'pie': XL_CHART_TYPE.PIE,
        'area': XL_CHART_TYPE.AREA,
        'scatter': XL_CHART_TYPE.XY_SCATTER,
        'bubble': XL_CHART_TYPE.BUBBLE,
        'doughnut': XL_CHART_TYPE.DOUGHNUT,
        'radar': XL_CHART_TYPE.RADAR,
        'bar3d': XL_CHART_TYPE.THREE_D_COLUMN_CLUSTERED,
    }
    return chart_type_map.get(chart_type_str.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)

def get_shape_type(shape_type_str: str) -> int:
    """Map shape type string to python-pptx shape type."""
    shape_type_map = {
        'rectangle': MSO_SHAPE.RECTANGLE,
        'ellipse': MSO_SHAPE.OVAL,
        'roundrectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
        'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'diamond': MSO_SHAPE.DIAMOND,
        'pentagon': MSO_SHAPE.REGULAR_PENTAGON,
        'hexagon': MSO_SHAPE.HEXAGON,
        'octagon': MSO_SHAPE.OCTAGON,
        'star': MSO_SHAPE.STAR_5_POINT,
        'arrow': MSO_SHAPE.RIGHT_ARROW,
    }
    return shape_type_map.get(shape_type_str.lower(), MSO_SHAPE.RECTANGLE)

async def generate_pptx(
    slides_data: List[Dict[str, Any]],
    options: Optional[Dict[str, Any]] = None
) -> bytes:
    """Generate PPTX from slide data using python-pptx."""
    if not slides_data:
        raise ValueError("At least one slide must be provided")
    
    # Set default options
    options = options or {}
    
    # Create new PowerPoint presentation
    print("🎯 Creating PowerPoint presentation...")
    prs = Presentation()
    
    # Set presentation properties (metadata)
    if options.get('author'):
        prs.core_properties.author = options['author']
    if options.get('title'):
        prs.core_properties.title = options['title']
    if options.get('subject'):
        prs.core_properties.subject = options['subject']
    
    # Set slide dimensions based on layout
    layout = options.get('layout', '16x9')
    if layout == '16x10':
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(8.5)
    elif layout == '4x3':
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:  # Default to 16x9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
    
    # Process each slide
    for i, slide_data in enumerate(slides_data):
        print(f"📄 Processing slide {i + 1} of {len(slides_data)}...")
        
        # Add slide with blank layout
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide background if specified
        if slide_data.get('backgroundColor'):
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(slide_data['backgroundColor'])
        
        # Process elements on the slide
        elements = slide_data.get('elements', [])
        for element in elements:
            element_type = element.get('type')
            element_options = element.get('options', {})
            
            try:
                if element_type == 'text':
                    await add_text_element(slide, element, element_options)
                elif element_type == 'table':
                    await add_table_element(slide, element, element_options)
                elif element_type == 'chart':
                    await add_chart_element(slide, element, element_options)
                elif element_type == 'shape':
                    await add_shape_element(slide, element, element_options)
                else:
                    print(f"❓ Unknown element type: {element_type}")
                    
            except Exception as element_error:
                print(f"❌ Error adding {element_type} element: {element_error}")
                # Continue processing other elements
    
    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

async def add_text_element(slide, element: Dict[str, Any], options: Dict[str, Any]):
    """Add a text element to the slide."""
    text_content = element.get('text', '')
    if not text_content:
        return
    
    # Set default position and size
    left = Inches(options.get('x', 1))
    top = Inches(options.get('y', 1))
    width = Inches(options.get('w', 8))
    height = Inches(options.get('h', 1))
    
    # Add text box
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    
    # Handle array of text objects or simple string
    if isinstance(text_content, list):
        for i, text_obj in enumerate(text_content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            if isinstance(text_obj, dict):
                run = p.add_run()
                run.text = text_obj.get('text', '')
                
                # Apply formatting
                if text_obj.get('bold'):
                    run.font.bold = True
                if text_obj.get('italic'):
                    run.font.italic = True
                if text_obj.get('fontSize'):
                    run.font.size = Pt(text_obj['fontSize'])
                if text_obj.get('color'):
                    run.font.color.rgb = hex_to_rgb(text_obj['color'])
            else:
                p.text = str(text_obj)
    else:
        text_frame.text = str(text_content)
        
        # Apply global formatting options
        if options.get('fontSize'):
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(options['fontSize'])
        
        if options.get('color'):
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = hex_to_rgb(options['color'])
        
        if options.get('bold'):
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

async def add_table_element(slide, element: Dict[str, Any], options: Dict[str, Any]):
    """Add a table element to the slide."""
    rows_data = element.get('rows', [])
    if not rows_data:
        return
    
    # Set default position and size
    left = Inches(options.get('x', 1))
    top = Inches(options.get('y', 1))
    width = Inches(options.get('w', 8))
    height = Inches(options.get('h', 3))
    
    # Create table
    rows = len(rows_data)
    cols = len(rows_data[0]) if rows_data else 1
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Populate table with data
    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_data in enumerate(row_data):
            if col_idx < len(table.rows[row_idx].cells):
                cell = table.rows[row_idx].cells[col_idx]
                
                if isinstance(cell_data, dict):
                    cell.text = str(cell_data.get('text', ''))
                else:
                    cell.text = str(cell_data)

async def add_chart_element(slide, element: Dict[str, Any], options: Dict[str, Any]):
    """Add a chart element to the slide."""
    chart_type = element.get('chartType')
    chart_data = element.get('chartData', [])
    
    if not chart_type or not chart_data:
        return
    
    # Set default position and size
    left = Inches(options.get('x', 1))
    top = Inches(options.get('y', 1))
    width = Inches(options.get('w', 8))
    height = Inches(options.get('h', 5))
    
    # Create chart data
    chart_data_obj = CategoryChartData()
    
    # Extract categories and series from chart data
    if chart_data and len(chart_data) > 0:
        first_series = chart_data[0]
        categories = first_series.get('labels', [])
        chart_data_obj.categories = categories
        
        for series in chart_data:
            series_name = series.get('name', 'Series')
            series_values = series.get('values', [])
            chart_data_obj.add_series(series_name, series_values)
        
        # Add chart to slide
        chart_type_enum = get_chart_type(chart_type)
        chart = slide.shapes.add_chart(chart_type_enum, left, top, width, height, chart_data_obj).chart
        
        # Set chart title if provided
        if options.get('title'):
            chart.chart_title.text_frame.text = options['title']

async def add_shape_element(slide, element: Dict[str, Any], options: Dict[str, Any]):
    """Add a shape element to the slide."""
    shape_type = element.get('shapeType')
    if not shape_type:
        return
    
    # Set default position and size
    left = Inches(options.get('x', 1))
    top = Inches(options.get('y', 1))
    width = Inches(options.get('w', 2))
    height = Inches(options.get('h', 2))
    
    # Add shape
    shape_type_enum = get_shape_type(shape_type)
    shape = slide.shapes.add_shape(shape_type_enum, left, top, width, height)
    
    # Apply formatting
    if options.get('fill'):
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(options['fill'])
    
    if options.get('line'):
        shape.line.color.rgb = hex_to_rgb(options['line'])

def get_file_size_string(content: bytes) -> str:
    """Calculate file size string from bytes content."""
    bytes_size = len(content)
    kb = bytes_size / 1024
    
    if kb < 1024:
        return f"{kb:.0f} KB" if kb >= 1 else "1 KB"
    else:
        return f"{kb / 1024:.2f} MB"

async def ensure_export_directory() -> None:
    """Ensure export directory exists, create if it doesn't."""
    export_path = Path(EXPORT_DIR)
    
    if export_path.exists():
        print(f"✓ Export directory exists: {EXPORT_DIR}")
    else:
        try:
            export_path.mkdir(parents=True, exist_ok=True)
            print(f"✓ Created export directory: {EXPORT_DIR}")
        except Exception as e:
            print(f"✗ Failed to create export directory: {e}")
            raise

async def write_pptx_to_file(pptx_content: bytes, filename: str) -> str:
    """Write PPTX content to file system."""
    await ensure_export_directory()
    
    filepath = Path(EXPORT_DIR) / filename
    
    try:
        filepath.write_bytes(pptx_content)
        print(f"✓ File written: {filepath}")
        return str(filepath)
    except Exception as e:
        print(f"✗ Failed to write file: {e}")
        raise

async def test_pptx_export():
    """Test the PPTX export functionality."""
    print("🧪 Testing PPTX Export Logic")
    print("=" * 50)
    
    try:
        # Extract arguments
        slides = TEST_ARGUMENTS.get("slides")
        filename = TEST_ARGUMENTS.get("filename", "output")
        description = TEST_ARGUMENTS.get("description")
        options = TEST_ARGUMENTS.get("options", {})
        
        print(f"📊 Slide Count: {len(slides)}")
        print(f"📄 Filename: {filename}")
        print(f"📝 Description: {description}")
        print(f"⚙️  Options: {options}")
        print()
        
        # Validate input
        if not slides or not isinstance(slides, list):
            raise ValueError("Slides must be provided as an array of slide objects")
        
        if len(slides) == 0:
            raise ValueError("At least one slide must be provided")
        
        # Generate PPTX
        print("🔄 Converting slide data to PPTX format...")
        pptx_content = await generate_pptx(slides, options)
        
        # Generate UUID and filename
        file_uuid = str(uuid.uuid4())
        sanitized_filename = "".join(c if c.isalnum() or c in "_-" else "_" for c in filename)
        full_filename = f"{sanitized_filename}_{file_uuid}.pptx"
        file_size = get_file_size_string(pptx_content)
        
        # Write PPTX to file system
        print("💾 Writing file to disk...")
        filepath = await write_pptx_to_file(pptx_content, full_filename)
        
        layout = options.get('layout', '16x9')
        print()
        print("✅ PPTX Export Successful!")
        print(f"📁 Generated file: {full_filename}")
        print(f"📏 File size: {file_size}")
        print(f"📊 Slides: {len(slides)}, Layout: {layout}")
        print(f"💾 Saved to: {filepath}")
        
        # Create result object (same as MCP server would return)
        result = {
            "path": full_filename,
            "filetype": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "filename": full_filename,
            "filesize": file_size,
        }
        
        print()
        print("📤 MCP Server Response:")
        print(json.dumps(result, indent=2))
        
    except Exception as error:
        print(f"❌ Error during PPTX export: {error}")
        
        error_result = {
            "success": False,
            "error": str(error),
        }
        
        print()
        print("📤 MCP Server Error Response:")
        print(json.dumps(error_result, indent=2))

if __name__ == "__main__":
    asyncio.run(test_pptx_export())