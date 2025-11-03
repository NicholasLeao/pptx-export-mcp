#!/usr/bin/env python3
"""PPTX Export MCP Server - Python implementation."""

import json
import sys
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional
import base64
import io

from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

# Export directory configuration
EXPORT_DIR = "/tmp/protex-intelligence-file-exports"


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor object."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def get_chart_type(chart_type_str: str) -> int:
    """Map chart type string to python-pptx chart type."""
    chart_type_map = {
        'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'line': XL_CHART_TYPE.LINE,
        'pie': XL_CHART_TYPE.PIE,
        'area': XL_CHART_TYPE.AREA,
        'scatter': XL_CHART_TYPE.XY_SCATTER,
        'bubble': XL_CHART_TYPE.BUBBLE,
        'doughnut': XL_CHART_TYPE.DOUGHNUT,
        'radar': XL_CHART_TYPE.RADAR,
        'bar3d': XL_CHART_TYPE.THREE_D_COLUMN_CLUSTERED,
    }
    return chart_type_map.get(chart_type_str.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)


def get_shape_type(shape_type_str: str) -> int:
    """Map shape type string to python-pptx shape type."""
    shape_type_map = {
        'rectangle': MSO_SHAPE.RECTANGLE,
        'ellipse': MSO_SHAPE.OVAL,
        'roundrectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
        'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'diamond': MSO_SHAPE.DIAMOND,
        'pentagon': MSO_SHAPE.REGULAR_PENTAGON,
        'hexagon': MSO_SHAPE.HEXAGON,
        'octagon': MSO_SHAPE.OCTAGON,
        'star': MSO_SHAPE.STAR_5_POINT,
        'arrow': MSO_SHAPE.RIGHT_ARROW,
    }
    return shape_type_map.get(shape_type_str.lower(), MSO_SHAPE.RECTANGLE)


async def generate_pptx(
    slides_data: List[Dict[str, Any]],
    options: Optional[Dict[str, Any]] = None
) -> bytes:
    """Generate PPTX from slide data using python-pptx."""
    if not slides_data:
        raise ValueError("At least one slide must be provided")
    
    # Set default options
    options = options or {}
    
    # Create new PowerPoint presentation
    print("🎯 Creating PowerPoint presentation...", file=sys.stderr)
    prs = Presentation()
    
    # Set presentation properties (metadata)
    if options.get('author'):
        prs.core_properties.author = options['author']
    if options.get('title'):
        prs.core_properties.title = options['title']
    if options.get('subject'):
        prs.core_properties.subject = options['subject']
    
    # Set slide dimensions based on layout
    layout = options.get('layout', '16x9')
    if layout == '16x10':
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(8.5)
    elif layout == '4x3':
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:  # Default to 16x9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
    
    # Process each slide
    for i, slide_data in enumerate(slides_data):
        print(f"📄 Processing slide {i + 1} of {len(slides_data)}...", file=sys.stderr)
        
        # Add slide with blank layout
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide background if specified
        if slide_data.get('backgroundColor'):
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(slide_data['backgroundColor'])
        
        # Process elements on the slide
        elements = slide_data.get('elements', [])
        for element in elements:
            element_type = element.get('type')
            element_options = element.get('options', {})
            
            try:
                if element_type == 'text':
                    await add_text_element(slide, element, element_options)
                elif element_type == 'table':
                    await add_table_element(slide, element, element_options)
                elif element_type == 'chart':
                    await add_chart_element(slide, element, element_options)
                elif element_type == 'image':
                    await add_image_element(slide, element, element_options)
                elif element_type == 'shape':
                    await add_shape_element(slide, element, element_options)
                else:
                    print(f"❓ Unknown element type: {element_type}", file=sys.stderr)
                    
            except Exception as element_error:
                print(f"❌ Error adding {element_type} element: {element_error}", file=sys.stderr)
                # Continue processing other elements
        
        # Add watermark to every slide
        try:
            await add_watermark(slide, prs)
        except Exception as watermark_error:
            print(f"❌ Error adding watermark: {watermark_error}", file=sys.stderr)
    
    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


async def add_text_element(slide, element: Dict[str, Any], options: Dict[str, Any]):
    """Add a text element to the slide."""
    text_content = element.get('text', '')
    if not text_content:
        return
    
    # Set default position and size
    left = Inches(options.get('x', 1))
    top = Inches(options.get('y', 1))
    width = Inches(options.get('w', 8))
    height = Inches(options.get('h', 1))
    
    # Add text box
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    
    # Handle array of text objects or simple string
    if isinstance(text_content, list):
        for i, text_obj in enumerate(text_content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            if isinstance(text_obj, dict):
                run = p.add_run()
                run.text = text_obj.get('text', '')
                
                # Apply formatting
                if text_obj.get('bold'):
                    run.font.bold = True
                if text_obj.get('italic'):
                    run.font.italic = True
                if text_obj.get('fontSize'):
                    run.font.size = Pt(text_obj['fontSize'])
                if text_obj.get('color'):
                    run.font.color.rgb = hex_to_rgb(text_obj['color'])
            else:
                p.text = str(text_obj)
    else:
        text_frame.text = str(text_content)
        
        # Apply global formatting options
        if options.get('fontSize'):
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(options['fontSize'])
        
        if options.get('color'):
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = hex_to_rgb(options['color'])
        
        if options.get('bold'):
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True


async def add_table_element(slide, element: Dict[str, Any], options: Dict[str, Any]):
    """Add a table element to the slide."""
    rows_data = element.get('rows', [])
    if not rows_data:
        return
    
    # Set default position and size
    left = Inches(options.get('x', 1))
    top = Inches(options.get('y', 1))
    width = Inches(options.get('w', 8))
    height = Inches(options.get('h', 3))
    
    # Create table
    rows = len(rows_data)
    cols = len(rows_data[0]) if rows_data else 1
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Populate table with data
    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_data in enumerate(row_data):
            if col_idx < len(table.rows[row_idx].cells):
                cell = table.rows[row_idx].cells[col_idx]
                
                if isinstance(cell_data, dict):
                    cell.text = str(cell_data.get('text', ''))
                    # Could add cell-specific formatting here
                else:
                    cell.text = str(cell_data)


async def add_chart_element(slide, element: Dict[str, Any], options: Dict[str, Any]):
    """Add a chart element to the slide."""
    chart_type = element.get('chartType')
    chart_data = element.get('chartData', [])
    
    if not chart_type or not chart_data:
        return
    
    # Set default position and size
    left = Inches(options.get('x', 1))
    top = Inches(options.get('y', 1))
    width = Inches(options.get('w', 8))
    height = Inches(options.get('h', 5))
    
    # Create chart data
    chart_data_obj = CategoryChartData()
    
    # Extract categories and series from chart data
    if chart_data and len(chart_data) > 0:
        first_series = chart_data[0]
        categories = first_series.get('labels', [])
        chart_data_obj.categories = categories
        
        for series in chart_data:
            series_name = series.get('name', 'Series')
            series_values = series.get('values', [])
            chart_data_obj.add_series(series_name, series_values)
        
        # Add chart to slide
        chart_type_enum = get_chart_type(chart_type)
        chart = slide.shapes.add_chart(chart_type_enum, left, top, width, height, chart_data_obj).chart
        
        # Set chart title if provided
        if options.get('title'):
            chart.chart_title.text_frame.text = options['title']


async def add_image_element(slide, element: Dict[str, Any], options: Dict[str, Any]):
    """Add an image element to the slide."""
    image_path = element.get('path')
    if not image_path:
        return
    
    # Set default position and size
    left = Inches(options.get('x', 1))
    top = Inches(options.get('y', 1))
    width = Inches(options.get('w', 4)) if options.get('w') else None
    height = Inches(options.get('h', 3)) if options.get('h') else None
    
    try:
        # Handle base64 data URI
        if image_path.startswith('data:'):
            # Extract base64 data
            header, data = image_path.split(',', 1)
            image_data = base64.b64decode(data)
            image_stream = io.BytesIO(image_data)
            
            if width and height:
                slide.shapes.add_picture(image_stream, left, top, width, height)
            else:
                slide.shapes.add_picture(image_stream, left, top)
        else:
            # Handle file path
            if width and height:
                slide.shapes.add_picture(image_path, left, top, width, height)
            else:
                slide.shapes.add_picture(image_path, left, top)
                
    except Exception as img_error:
        print(f"❌ Error adding image: {img_error}", file=sys.stderr)


async def add_shape_element(slide, element: Dict[str, Any], options: Dict[str, Any]):
    """Add a shape element to the slide."""
    shape_type = element.get('shapeType')
    if not shape_type:
        return
    
    # Set default position and size
    left = Inches(options.get('x', 1))
    top = Inches(options.get('y', 1))
    width = Inches(options.get('w', 2))
    height = Inches(options.get('h', 2))
    
    # Add shape
    shape_type_enum = get_shape_type(shape_type)
    shape = slide.shapes.add_shape(shape_type_enum, left, top, width, height)
    
    # Apply formatting
    if options.get('fill'):
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(options['fill'])
    
    if options.get('line'):
        shape.line.color.rgb = hex_to_rgb(options['line'])


async def add_watermark(slide, prs):
    """Add Protex AI watermark to the bottom right of the slide."""
    # Calculate position based on slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Position watermark in bottom right with some padding
    watermark_width = Inches(2)
    watermark_height = Inches(0.3)
    left = slide_width - watermark_width - Inches(0.2)  # 0.2" padding from right
    top = slide_height - watermark_height - Inches(0.2)  # 0.2" padding from bottom
    
    # Add watermark text box
    watermark_textbox = slide.shapes.add_textbox(left, top, watermark_width, watermark_height)
    text_frame = watermark_textbox.text_frame
    text_frame.text = "Powered by Protex AI"
    
    # Style the watermark text
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.RIGHT
    
    for run in paragraph.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(102, 102, 102)  # Gray color (#666666)
        run.font.name = "Arial"


def get_file_size_string(content: bytes) -> str:
    """Calculate file size string from bytes content."""
    bytes_size = len(content)
    kb = bytes_size / 1024
    
    if kb < 1024:
        return f"{kb:.0f} KB" if kb >= 1 else "1 KB"
    else:
        return f"{kb / 1024:.2f} MB"


async def ensure_export_directory() -> None:
    """Ensure export directory exists, create if it doesn't."""
    export_path = Path(EXPORT_DIR)
    
    if export_path.exists():
        print(f"✓ Export directory exists: {EXPORT_DIR}", file=sys.stderr)
    else:
        try:
            export_path.mkdir(parents=True, exist_ok=True)
            print(f"✓ Created export directory: {EXPORT_DIR}", file=sys.stderr)
        except Exception as e:
            print(f"✗ Failed to create export directory: {e}", file=sys.stderr)
            raise


async def write_pptx_to_file(pptx_content: bytes, filename: str) -> str:
    """Write PPTX content to file system."""
    await ensure_export_directory()
    
    filepath = Path(EXPORT_DIR) / filename
    
    try:
        filepath.write_bytes(pptx_content)
        print(f"✓ File written: {filepath}", file=sys.stderr)
        return str(filepath)
    except Exception as e:
        print(f"✗ Failed to write file: {e}", file=sys.stderr)
        raise


# Create FastMCP server
mcp = FastMCP("pptx-export-mcp")


@mcp.tool()
async def pptx_export(
    slides: List[Dict[str, Any]],
    filename: str = "output",
    description: str = None,
    options: Dict[str, Any] = None
) -> dict:
    """Export data to PowerPoint (PPTX) format with full support for text, tables, charts, images, and shapes.
    
    Args:
        slides: Array of slide objects containing elements to add
        filename: Filename for the exported file (without extension)
        description: Optional description of the file contents
        options: Presentation options including layout, author, title, subject
        
    Returns:
        Dictionary with export results including path and file info
    """
    try:
        # Validate input
        if not slides or not isinstance(slides, list):
            raise ValueError("Slides must be provided as an array of slide objects")
        
        if len(slides) == 0:
            raise ValueError("At least one slide must be provided")
        
        # Generate PPTX
        print("🔄 Generating PowerPoint from slide data...", file=sys.stderr)
        pptx_content = await generate_pptx(slides, options or {})
        
        # Generate UUID and filename
        file_uuid = str(uuid.uuid4())
        sanitized_filename = "".join(c if c.isalnum() or c in "_-" else "_" for c in filename)
        full_filename = f"{sanitized_filename}_{file_uuid}.pptx"
        file_size = get_file_size_string(pptx_content)
        
        # Write PPTX to file system
        filepath = await write_pptx_to_file(pptx_content, full_filename)
        
        layout = (options or {}).get('layout', '16x9')
        print(f"✅ PPTX generated: {full_filename} ({file_size})", file=sys.stderr)
        print(f"   Slides: {len(slides)}, Layout: {layout}", file=sys.stderr)
        print(f"   Saved to: {filepath}", file=sys.stderr)
        
        # Return simplified response with essential information
        return {
            "path": full_filename,
            "filetype": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "filename": full_filename,
            "filesize": file_size,
        }
        
    except Exception as error:
        print(f"Error processing PPTX export: {error}", file=sys.stderr)
        
        return {
            "success": False,
            "error": str(error),
        }


def cli_main():
    """CLI entry point."""
    mcp.run()


if __name__ == "__main__":
    cli_main()